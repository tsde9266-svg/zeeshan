import os
import re
from bs4 import BeautifulSoup
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

def html_to_pptx_converter(html_file_path, output_pptx_path, image_folder='images'):
    """
    Convert the HTML diabetes presentation to PPTX format with improved styling
    """
    # Read and parse HTML
    with open(html_file_path, 'r', encoding='utf-8') as file:
        html_content = file.read()
    
    soup = BeautifulSoup(html_content, 'html.parser')
    
    # Create presentation
    prs = Presentation()
    prs.slide_width = Inches(13.33)  # Widescreen for better layout
    prs.slide_height = Inches(7.5)
    
    # Define slide layouts
    title_slide_layout = prs.slide_layouts[0]
    content_slide_layout = prs.slide_layouts[1]
    blank_slide_layout = prs.slide_layouts[5]
    
    # Extract all slides
    slides_data = extract_slides_data(soup)
    
    # Process each slide
    for i, slide_data in enumerate(slides_data):
        if i == 0:  # Title slide
            slide = prs.slides.add_slide(title_slide_layout)
            create_title_slide(slide, slide_data)
        else:  # Content slides
            slide = prs.slides.add_slide(blank_slide_layout)  # Use blank for custom layout
            create_content_slide(slide, slide_data, image_folder)
        
        # Add slide counter
        add_slide_counter(slide, i + 1, len(slides_data))
    
    # Save presentation
    prs.save(output_pptx_path)
    print(f"✅ PPTX saved as: {output_pptx_path}")
    print("📝 Note: Gradients, animations, and rounded corners may need manual adjustment in PowerPoint.")

def extract_slides_data(soup):
    """Extract structured data from each slide"""
    slides = soup.find_all('div', class_='slide')
    slides_data = []
    
    for slide in slides:
        slide_info = {
            'title': '',
            'content': [],
            'stats': [],
            'tables': [],
            'bullet_points': [],
            'highlight_boxes': [],
            'algorithm_grid': [],
            'feature_bars': [],
            'chart_placeholders': []
        }
        
        # Extract title
        title_elem = slide.find(['h1', 'h2'])
        if title_elem:
            slide_info['title'] = title_elem.get_text().strip()
        
        # Extract statistics cards
        stat_cards = slide.find_all('div', class_='stat-card')
        for card in stat_cards:
            number = card.find('span', class_='stat-number')
            label = card.find('span', class_='stat-label')
            if number and label:
                slide_info['stats'].append({
                    'number': number.get_text().strip(),
                    'label': label.get_text().strip()
                })
        
        # Extract bullet points
        bullet_lists = slide.find_all('ul', class_='bullet-points')
        for ul in bullet_lists:
            points = [li.get_text().strip() for li in ul.find_all('li')]
            slide_info['bullet_points'].extend(points)
        
        # Extract tables
        tables = slide.find_all('table', class_='performance-table')
        for table in tables:
            table_data = extract_table_data(table)
            slide_info['tables'].append(table_data)
        
        # Extract highlight boxes
        highlight_boxes = slide.find_all('div', class_='highlight-box')
        for box in highlight_boxes:
            text = box.get_text().strip()
            if text:
                slide_info['highlight_boxes'].append(text)
        
        # Extract algorithm grid
        algo_grids = slide.find_all('div', class_='algorithm-grid')
        for grid in algo_grids:
            cards = grid.find_all('div', class_='algorithm-card')
            for card in cards:
                title = card.find('h4').get_text().strip() if card.find('h4') else ''
                text = card.find('p').get_text().strip() if card.find('p') else ''
                slide_info['algorithm_grid'].append({'title': title, 'text': text})
        
        # Extract feature importance bars
        feature_bars = slide.find_all('div', class_='feature-bar')
        for bar in feature_bars:
            name = bar.find('span', class_='feature-name')
            value = bar.find('div', class_='bar-fill')
            if name and value:
                width_match = re.search(r'width:\s*(\d+)%', value.get('style', ''))
                width = int(width_match.group(1)) if width_match else 0
                slide_info['feature_bars'].append({
                    'name': name.get_text().strip(),
                    'value': value.get_text().strip(),
                    'width': width
                })
        
        # Extract chart placeholders
        chart_placeholders = slide.find_all('div', class_='chart-placeholder')
        for placeholder in chart_placeholders:
            slide_info['chart_placeholders'].append(placeholder.get_text().strip())
        
        # Extract regular paragraphs
        paragraphs = slide.find_all('p')
        for p in paragraphs:
            if not p.find_parent(['div'], class_=['stat-card', 'highlight-box', 'algorithm-card']):
                text = p.get_text().strip()
                if text:
                    slide_info['content'].append(text)
        
        slides_data.append(slide_info)
    
    return slides_data

def extract_table_data(table):
    """Extract table data as structured format"""
    headers = []
    rows = []
    
    thead = table.find('thead')
    if thead:
        header_row = thead.find('tr')
        headers = [th.get_text().strip() for th in header_row.find_all('th')]
    
    tbody = table.find('tbody') or table
    for tr in tbody.find_all('tr'):
        row_data = [td.get_text().strip() for td in tr.find_all(['td', 'th'])]
        if row_data:
            rows.append(row_data)
    
    return {'headers': headers, 'rows': rows}

def add_slide_counter(slide, current, total):
    """Add slide counter to top-right corner"""
    left = Inches(12)
    top = Inches(0.2)
    width = Inches(1)
    height = Inches(0.5)
    
    counter_shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top, width, height
    )
    counter_shape.fill.solid()
    counter_shape.fill.fore_color.rgb = RGBColor(0, 0, 0)
    counter_shape.fill.fore_color.brightness = -0.3  # Semi-transparent black
    counter_shape.line.color.rgb = RGBColor(0, 0, 0)
    
    text_frame = counter_shape.text_frame
    text_frame.clear()
    p = text_frame.paragraphs[0]
    p.text = f"{current} / {total}"
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER

def create_title_slide(slide, slide_data):
    """Create the title slide"""
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = slide_data['title']
    
    subtitle_text = "\n".join(slide_data['content'])
    subtitle.text = subtitle_text or """A Comparative Analysis Using the Pima Indian Diabetes Dataset\n\nPardeep Singh\nMSc. Computer Science and Technology\nUlster University, Birmingham, United Kingdom\n\nCW-1 Video Demonstration"""
    
    # Format title
    title_paragraph = title.text_frame.paragraphs[0]
    title_paragraph.font.size = Pt(36)
    title_paragraph.font.bold = True
    title_paragraph.font.color.rgb = RGBColor(44, 62, 80)  # #2c3e50
    title_paragraph.alignment = PP_ALIGN.CENTER
    
    # Format subtitle
    for paragraph in subtitle.text_frame.paragraphs:
        paragraph.font.size = Pt(18)
        paragraph.font.color.rgb = RGBColor(51, 51, 51)
        paragraph.alignment = PP_ALIGN.CENTER
    
    # Add highlight box if present
    if slide_data['highlight_boxes']:
        left = Inches(2)
        top = Inches(5)
        width = Inches(9)
        height = Inches(1.5)
        box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
        box.fill.solid()
        box.fill.fore_color.rgb = RGBColor(255, 243, 205)  # #fff3cd
        box.line.color.rgb = RGBColor(255, 234, 167)  # #ffeaa7
        
        text_frame = box.text_frame
        text_frame.clear()
        p = text_frame.paragraphs[0]
        p.text = slide_data['highlight_boxes'][0]
        p.font.size = Pt(14)
        p.font.color.rgb = RGBColor(51, 51, 51)
        p.alignment = PP_ALIGN.CENTER

def create_content_slide(slide, slide_data, image_folder):
    """Create content slides with improved styling"""
    # Add title
    title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12), Inches(1))
    title_tf = title_shape.text_frame
    title_tf.clear()
    p = title_tf.paragraphs[0]
    p.text = slide_data['title']
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RGBColor(52, 73, 94)  # #34495e
    p.alignment = PP_ALIGN.LEFT
    
    # Starting position for content
    top = Inches(1.8)
    left = Inches(0.5)
    
    # Add chart placeholders as images
    if slide_data['chart_placeholders']:
        for placeholder in slide_data['chart_placeholders']:
            img_name = re.sub(r'[\[\]]', '', placeholder).replace(' ', '_').replace('-', '_') + '.png'
            img_path = os.path.join(image_folder, img_name)
            if os.path.exists(img_path):
                pic = slide.shapes.add_picture(img_path, Inches(2), top, width=Inches(9))
                top += Inches(3.5)
            else:
                print(f"⚠️ Image not found: {img_path}")
    
    # Add statistics
    if slide_data['stats']:
        top = add_statistics_to_slide(slide, slide_data['stats'], top)
    
    # Add algorithm grid
    if slide_data['algorithm_grid']:
        top = add_algorithm_grid_to_slide(slide, slide_data['algorithm_grid'], top)
    
    # Add feature importance bars
    if slide_data['feature_bars']:
        top = add_feature_bars_to_slide(slide, slide_data['feature_bars'], top)
    
    # Add bullet points
    if slide_data['bullet_points']:
        textbox = slide.shapes.add_textbox(left, top, Inches(12), Inches(3))
        text_frame = textbox.text_frame
        text_frame.clear()
        for i, point in enumerate(slide_data['bullet_points'][:8]):
            p = text_frame.paragraphs[0] if i == 0 else text_frame.add_paragraph()
            p.text = point
            p.level = 0
            p.font.size = Pt(16)
            p.font.color.rgb = RGBColor(51, 51, 51)
            p.alignment = PP_ALIGN.LEFT
            run = p.runs[0] if p.runs else p.add_run()
            run.text = f"▶ {p.text}"
            run.font.color.rgb = RGBColor(52, 152, 219)  # #3498db for bullet
            p.text = p.text[2:] if p.text.startswith('▶ ') else p.text
        top += Inches(3.5)
    
    # Add tables
    if slide_data['tables']:
        top = add_table_to_slide(slide, slide_data['tables'][0], top)
    
    # Add highlight boxes
    if slide_data['highlight_boxes']:
        for box_text in slide_data['highlight_boxes'][:2]:  # Limit to 2 boxes
            box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, Inches(12), Inches(1))
            box.fill.solid()
            box.fill.fore_color.rgb = RGBColor(255, 243, 205)  # #fff3cd
            box.line.color.rgb = RGBColor(255, 234, 167)  # #ffeaa7
            text_frame = box.text_frame
            text_frame.clear()
            p = text_frame.paragraphs[0]
            p.text = box_text
            p.font.size = Pt(14)
            p.font.color.rgb = RGBColor(51, 51, 51)
            p.alignment = PP_ALIGN.LEFT
            top += Inches(1.2)

def add_statistics_to_slide(slide, stats, top):
    """Add statistics cards to slide"""
    stats_per_row = min(4, len(stats))
    card_width = Inches(2.8)
    card_height = Inches(1.8)
    left = Inches(0.5)
    
    for i, stat in enumerate(stats[:8]):  # Limit to 8 stats
        row = i // stats_per_row
        col = i % stats_per_row
        x = left + col * (card_width + Inches(0.3))
        y = top + row * (card_height + Inches(0.3))
        
        stat_shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, x, y, card_width, card_height
        )
        stat_shape.fill.solid()
        stat_shape.fill.fore_color.rgb = RGBColor(52, 152, 219)  # #3498db
        stat_shape.line.color.rgb = RGBColor(41, 128, 185)  # #2980b9
        stat_shape.line.width = Pt(1)
        
        text_frame = stat_shape.text_frame
        text_frame.clear()
        number_p = text_frame.paragraphs[0]
        number_p.text = stat['number']
        number_p.font.size = Pt(28)
        number_p.font.bold = True
        number_p.font.color.rgb = RGBColor(255, 255, 255)
        number_p.alignment = PP_ALIGN.CENTER
        
        label_p = text_frame.add_paragraph()
        label_p.text = stat['label']
        label_p.font.size = Pt(12)
        label_p.font.color.rgb = RGBColor(255, 255, 255)
        label_p.alignment = PP_ALIGN.CENTER
    
    return top + (len(stats) // stats_per_row + 1) * (card_height + Inches(0.3))

def add_algorithm_grid_to_slide(slide, algo_cards, top):
    """Add algorithm grid to slide"""
    cards_per_row = 2
    card_width = Inches(5.5)
    card_height = Inches(2)
    left = Inches(0.5)
    
    for i, card in enumerate(algo_cards[:4]):  # Limit to 4 cards
        row = i // cards_per_row
        col = i % cards_per_row
        x = left + col * (card_width + Inches(0.3))
        y = top + row * (card_height + Inches(0.3))
        
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, x, y, card_width, card_height
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(248, 249, 250)  # #f8f9fa
        shape.line.color.rgb = RGBColor(23, 162, 184)  # #17a2b8
        shape.line.width = Pt(3)
        
        text_frame = shape.text_frame
        text_frame.clear()
        p = text_frame.paragraphs[0]
        p.text = card['title']
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = RGBColor(51, 51, 51)
        p.alignment = PP_ALIGN.LEFT
        
        p2 = text_frame.add_paragraph()
        p2.text = card['text']
        p2.font.size = Pt(12)
        p2.font.color.rgb = RGBColor(51, 51, 51)
        p2.alignment = PP_ALIGN.LEFT
    
    return top + (len(algo_cards) // cards_per_row + 1) * (card_height + Inches(0.3))

def add_feature_bars_to_slide(slide, feature_bars, top):
    """Add feature importance bars to slide"""
    bar_width = Inches(9)
    bar_height = Inches(0.5)
    left = Inches(2)
    
    for i, bar in enumerate(feature_bars[:5]):  # Limit to 5 bars
        y = top + i * (bar_height + Inches(0.2))
        
        # Add feature name
        name_shape = slide.shapes.add_textbox(Inches(0.5), y, Inches(1.5), bar_height)
        name_tf = name_shape.text_frame
        name_tf.clear()
        p = name_tf.paragraphs[0]
        p.text = bar['name']
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = RGBColor(51, 51, 51)
        
        # Add bar container
        container = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, left, y, bar_width, bar_height
        )
        container.fill.solid()
        container.fill.fore_color.rgb = RGBColor(236, 240, 241)  # #ecf0f1
        container.line.color.rgb = RGBColor(189, 195, 199)  # #bdc3c7
        
        # Add bar fill
        fill_width = Inches(bar_width.inches * (bar['width'] / 100))
        fill = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, left, y, fill_width, bar_height
        )
        fill.fill.solid()
        fill.fill.fore_color.rgb = RGBColor(231, 76, 60)  # #e74c3c
        fill.line.color.rgb = RGBColor(192, 57, 43)  # #c0392b
        
        # Add value text
        text_shape = slide.shapes.add_textbox(
            left + fill_width - Inches(0.5), y, Inches(0.5), bar_height
        )
        text_tf = text_shape.text_frame
        text_tf.clear()
        p = text_tf.paragraphs[0]
        p.text = bar['value']
        p.font.size = Pt(12)
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.alignment = PP_ALIGN.RIGHT
    
    return top + (len(feature_bars) * (bar_height + Inches(0.2)))

def add_table_to_slide(slide, table_data, top):
    """Add table to slide with improved styling"""
    if not table_data['headers'] or not table_data['rows']:
        return top
    
    rows = len(table_data['rows']) + 1
    cols = len(table_data['headers'])
    left = Inches(0.5)
    width = Inches(12)
    height = Inches(min(4, rows * 0.5))
    
    table = slide.shapes.add_table(rows, cols, left, top, width, height).table
    
    for col_idx in range(cols):
        table.columns[col_idx].width = Inches(width.inches / cols)
    
    for col_idx, header in enumerate(table_data['headers']):
        cell = table.cell(0, col_idx)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(44, 62, 80)  # #2c3e50
        paragraph = cell.text_frame.paragraphs[0]
        paragraph.font.color.rgb = RGBColor(255, 255, 255)
        paragraph.font.bold = True
        paragraph.font.size = Pt(12)
        paragraph.alignment = PP_ALIGN.LEFT
    
    for row_idx, row_data in enumerate(table_data['rows'][:10]):
        for col_idx, cell_data in enumerate(row_data[:cols]):
            if col_idx < cols:
                cell = table.cell(row_idx + 1, col_idx)
                cell.text = str(cell_data)
                cell.fill.solid()
                if 'best-score' in str(row_data).lower():
                    cell.fill.fore_color.rgb = RGBColor(212, 237, 218)  # #d4edda
                elif row_idx % 2 == 0:
                    cell.fill.fore_color.rgb = RGBColor(248, 249, 250)  # #f8f9fa
                else:
                    cell.fill.fore_color.rgb = RGBColor(255, 255, 255)
                paragraph = cell.text_frame.paragraphs[0]
                paragraph.font.size = Pt(10)
                paragraph.font.color.rgb = RGBColor(51, 51, 51)
                paragraph.alignment = PP_ALIGN.LEFT
    
    return top + height + Inches(0.3)

def main():
    """Main function to run the converter"""
    print("🔄 HTML to PPTX Converter for Diabetes Presentation")
    print("=" * 50)
    
    html_file = "index.html"
    output_file = "Machine_Learning_Diabetes_Prediction.pptx"
    image_folder = "images"
    
    if not os.path.exists(html_file):
        print(f"❌ HTML file not found: {html_file}")
        return
    
    if not os.path.exists(image_folder):
        print(f"⚠️ Image folder not found: {image_folder}. Creating it.")
        os.makedirs(image_folder)
    
    try:
        html_to_pptx_converter(html_file, output_file, image_folder)
        print(f"✅ Successfully converted to: {output_file}")
    except Exception as e:
        print(f"❌ Error during conversion: {str(e)}")

if __name__ == "__main__":
    print("Installing required packages...")
    os.system("pip install python-pptx beautifulsoup4 lxml")
    main()